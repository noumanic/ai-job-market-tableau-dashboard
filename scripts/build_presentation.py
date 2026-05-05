"""
Build the Group 8 final presentation as a Canva-style PPTX.

Output: docs/Group8_Final_Presentation.pptx
"""
from pathlib import Path
from copy import deepcopy
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "docs" / "figures"
ASSETS = ROOT / "assets"
OUT = ROOT / "docs" / "Group8_Final_Presentation.pptx"

# ---- Forest Emerald palette ---------------------------------------------
FOREST       = RGBColor(0x1B, 0x43, 0x32)
EMERALD      = RGBColor(0x2D, 0x6A, 0x4F)
JADE         = RGBColor(0x40, 0x91, 0x6C)
SAGE         = RGBColor(0x52, 0xB7, 0x88)
MINT         = RGBColor(0xB7, 0xE4, 0xC7)
HONEY        = RGBColor(0xD8, 0x97, 0x3C)
RUST         = RGBColor(0xBC, 0x47, 0x49)
CREAM        = RGBColor(0xFA, 0xF8, 0xF0)
LIGHT_MINT   = RGBColor(0xE8, 0xF1, 0xED)
PEACH        = RGBColor(0xFF, 0xF4, 0xE5)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK         = RGBColor(0x1A, 0x1A, 0x1A)
GRAY         = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]


# ---- helpers ------------------------------------------------------------
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0, shape_type=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape_type, x, y, w, h)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w if line_w else 1)
    return sh


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
             color=DARK, align=PP_ALIGN.LEFT, font="Calibri",
             anchor=MSO_ANCHOR.TOP, line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tb


def add_multi_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT,
                   anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """runs is a list of dicts: {text, size, bold, color, italic, font}"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    for spec in runs:
        if spec.get("newline"):
            p = tf.add_paragraph()
            p.alignment = spec.get("align", align)
            if line_spacing:
                p.line_spacing = line_spacing
            continue
        r = p.add_run()
        r.text = spec["text"]
        r.font.size = Pt(spec.get("size", 14))
        r.font.bold = spec.get("bold", False)
        r.font.italic = spec.get("italic", False)
        r.font.color.rgb = spec.get("color", DARK)
        r.font.name = spec.get("font", "Calibri")
    return tb


def slide_chrome(title, subtitle=None, page_n=None, page_total=28):
    slide = prs.slides.add_slide(BLANK)

    # Page background — soft cream
    add_rect(slide, 0, 0, SW, SH, fill=CREAM)

    # Top accent bar (forest)
    add_rect(slide, 0, 0, SW, Inches(0.22), fill=FOREST)
    # Sub accent strip (sage)
    add_rect(slide, 0, Inches(0.22), Inches(2.0), Inches(0.06), fill=SAGE)

    # Course tag (top left, on the strip)
    add_text(slide, Inches(0.4), Inches(0.0), Inches(8), Inches(0.22),
             "AI JOB MARKET INTELLIGENCE  ·  CS3012  ·  GROUP 8",
             size=9, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

    # Page number (top right)
    if page_n is not None:
        add_text(slide, SW - Inches(1.5), Inches(0.0), Inches(1.2), Inches(0.22),
                 f"{page_n:02d} / {page_total:02d}",
                 size=9, bold=True, color=WHITE, align=PP_ALIGN.RIGHT,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Title
    add_text(slide, Inches(0.5), Inches(0.45), SW - Inches(1.0), Inches(0.7),
             title, size=32, bold=True, color=FOREST,
             anchor=MSO_ANCHOR.MIDDLE, font="Calibri")

    # Subtitle (optional)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.15), SW - Inches(1.0), Inches(0.4),
                 subtitle, size=14, italic=True, color=GRAY,
                 anchor=MSO_ANCHOR.TOP)

    # Bottom footer strip (forest)
    add_rect(slide, 0, SH - Inches(0.25), SW, Inches(0.25), fill=FOREST)
    add_text(slide, Inches(0.4), SH - Inches(0.25), SW - Inches(0.8), Inches(0.25),
             "FAST-NUCES Islamabad  ·  Spring 2026  ·  Dr. Atif Mughees",
             size=9, color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(0.4), SH - Inches(0.25), SW - Inches(0.8), Inches(0.25),
             "Source: aijobs.net (CC0)",
             size=9, color=WHITE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    return slide


def card(slide, x, y, w, h, fill=WHITE, border=EMERALD, radius=True):
    """Decorative rounded card with sage border."""
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.shadow.inherit = False
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = border
    sh.line.width = Pt(2)
    return sh


def kpi_tile(slide, x, y, w, h, value, label, accent):
    card(slide, x, y, w, h, fill=CREAM, border=accent)
    add_text(slide, x, y + Inches(0.15), w, Inches(0.7),
             value, size=30, bold=True, color=accent,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, y + h - Inches(0.4), w, Inches(0.3),
             label, size=10, bold=True, color=GRAY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def icon_circle(slide, x, y, diameter, color, icon_text, font_size=18):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, diameter, diameter)
    sh.shadow.inherit = False
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    add_text(slide, x, y, diameter, diameter,
             icon_text, size=font_size, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# =========================================================================
# SLIDE 1 — TITLE
# =========================================================================
def slide_title():
    slide = prs.slides.add_slide(BLANK)

    # Full forest background
    add_rect(slide, 0, 0, SW, SH, fill=FOREST)

    # Diagonal sage accent (decorative band)
    add_rect(slide, Inches(-1.5), Inches(5.0), Inches(8), Inches(2.0), fill=EMERALD)
    add_rect(slide, Inches(7.0), Inches(0.0), Inches(8), Inches(1.5), fill=EMERALD)

    # Sage corner triangles
    tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                 Inches(11.0), Inches(5.5), Inches(2.5), Inches(2.0))
    tri.shadow.inherit = False
    tri.fill.solid(); tri.fill.fore_color.rgb = SAGE
    tri.line.fill.background()

    # Sage horizontal accent line
    add_rect(slide, Inches(0.7), Inches(2.6), Inches(1.5), Inches(0.06), fill=SAGE)

    # Course header
    add_text(slide, Inches(0.7), Inches(2.0), Inches(10), Inches(0.4),
             "CS3012  ·  FUNDAMENTALS OF DATA VISUALIZATION  ·  SPRING 2026",
             size=12, bold=True, color=SAGE, font="Calibri")

    # Title
    add_text(slide, Inches(0.7), Inches(2.8), Inches(12), Inches(1.2),
             "AI JOB MARKET",
             size=72, bold=True, color=WHITE, font="Calibri")
    add_text(slide, Inches(0.7), Inches(3.8), Inches(12), Inches(1.2),
             "INTELLIGENCE",
             size=72, bold=True, color=SAGE, font="Calibri")

    # Subtitle
    add_text(slide, Inches(0.7), Inches(5.0), Inches(11), Inches(0.5),
             "Salaries  ·  Roles  ·  Global Hiring Trends  ·  2022 – 2025",
             size=20, italic=True, color=CREAM, font="Calibri")

    # Stats row
    add_multi_text(slide, Inches(0.7), Inches(5.6), Inches(11), Inches(0.4), runs=[
        {"text": "147,348", "size": 14, "bold": True, "color": SAGE},
        {"text": "  records  ·  ",   "size": 14, "color": CREAM},
        {"text": "87",      "size": 14, "bold": True, "color": SAGE},
        {"text": "  countries  ·  ", "size": 14, "color": CREAM},
        {"text": "15",      "size": 14, "bold": True, "color": SAGE},
        {"text": "  role categories  ·  ", "size": 14, "color": CREAM},
        {"text": "10",      "size": 14, "bold": True, "color": SAGE},
        {"text": "  Tableau worksheets",   "size": 14, "color": CREAM},
    ])

    # Bottom-right: GROUP 8
    add_text(slide, Inches(9.5), Inches(6.4), Inches(3.5), Inches(0.4),
             "GROUP 8", size=14, bold=True, color=WHITE,
             align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(9.5), Inches(6.7), Inches(3.5), Inches(0.4),
             "FAST-NUCES Islamabad", size=11, color=CREAM,
             align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(9.5), Inches(6.95), Inches(3.5), Inches(0.4),
             "Instructor: Dr. Atif Mughees", size=11, color=CREAM,
             align=PP_ALIGN.RIGHT)

    return slide


# =========================================================================
# SLIDE 2 — AGENDA
# =========================================================================
def slide_agenda():
    slide = slide_chrome("Agenda", "What we will cover today", page_n=2)

    items = [
        ("01", "Group & Work Division",        "Who built what",                      FOREST),
        ("02", "Project Overview",             "Goals and the 5 guiding questions",   EMERALD),
        ("03", "About the Dataset",            "Source, properties, authenticity",    JADE),
        ("04", "Preprocessing Pipeline",       "5 stages from 151,445 to 147,348",    SAGE),
        ("05", "Exploratory Data Analysis",    "20 figures across the market",        HONEY),
        ("06", "Tableau Master Dashboard",     "10 worksheets · 5 KPIs · 4 insights", RUST),
        ("07", "Story the Dashboard Tells",    "Narrative that the visuals deliver",  EMERALD),
        ("08", "Key Insights",                 "What the data revealed",              FOREST),
        ("09", "How Tableau Helped",           "From static figures to a tool",       JADE),
        ("10", "Conclusion",                   "Wrap up & takeaways",                 SAGE),
    ]

    cols = 2
    rows = 5
    cw, ch = Inches(6.0), Inches(0.95)
    x0, y0 = Inches(0.65), Inches(1.7)
    gap_x, gap_y = Inches(0.2), Inches(0.15)

    for i, (num, title, sub, color) in enumerate(items):
        col = i % cols
        row = i // cols
        x = x0 + col * (cw + gap_x)
        y = y0 + row * (ch + gap_y)

        card(slide, x, y, cw, ch, fill=WHITE, border=color)
        # Big number on left
        add_text(slide, x + Inches(0.15), y + Inches(0.05), Inches(1.0), ch - Inches(0.1),
                 num, size=32, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
        # Title + sub
        add_text(slide, x + Inches(1.2), y + Inches(0.10), cw - Inches(1.4), Inches(0.45),
                 title, size=15, bold=True, color=FOREST)
        add_text(slide, x + Inches(1.2), y + Inches(0.45), cw - Inches(1.4), Inches(0.45),
                 sub, size=10, color=GRAY, italic=True)


# =========================================================================
# SLIDE 3 — GROUP MEMBERS & WORK DIVISION
# =========================================================================
def slide_group():
    slide = slide_chrome("Group 8 — Team & Work Division", "Three contributors, one dashboard", page_n=3)

    members = [
        ("Muhammad Nouman Hafeez", "21i-0416", "Project Lead & Tableau Engineering",
         ["Designed master dashboard layout and forest emerald palette",
          "Built all 10 Tableau worksheets and KPI tiles",
          "Authored Final_Project_Report.tex and dashboard build plan",
          "Wired filter actions, calculated fields, sorting"], FOREST),
        ("Muhammad Asim",          "21i-0852", "Data Engineering & Preprocessing",
         ["Wrote 5-script Python preprocessing pipeline",
          "Cleaned 151,445 → 147,348 rows; engineered 12 columns",
          "Mapped 406 raw job titles to 15 standardized categories",
          "Produced summary CSVs and dataset documentation"], EMERALD),
        ("Sara Jabeen",            "21i-0624", "EDA & Visual Story",
         ["Built EDA notebook, generated 20 matplotlib/seaborn figures",
          "Authored narrative analysis and 8 insight callouts",
          "Verified data authenticity with 4 source-platform screenshots",
          "Co-wrote final report and presentation deck"], JADE),
    ]

    cw = Inches(4.05); ch = Inches(5.0)
    x0 = Inches(0.45); y0 = Inches(1.7)
    gap = Inches(0.2)

    for i, (name, roll, role, items, color) in enumerate(members):
        x = x0 + i * (cw + gap)
        # card
        card(slide, x, y0, cw, ch, fill=WHITE, border=color)
        # color top strip
        add_rect(slide, x, y0, cw, Inches(0.5), fill=color, line=None)
        # icon circle
        icon_circle(slide, x + cw/2 - Inches(0.45), y0 + Inches(0.6),
                    Inches(0.9), color, ["MN","MA","SJ"][i], font_size=22)
        # name
        add_text(slide, x, y0 + Inches(1.6), cw, Inches(0.4),
                 name, size=15, bold=True, color=FOREST,
                 align=PP_ALIGN.CENTER)
        # roll
        add_text(slide, x, y0 + Inches(1.95), cw, Inches(0.3),
                 roll, size=11, color=GRAY, align=PP_ALIGN.CENTER)
        # role
        add_text(slide, x, y0 + Inches(2.25), cw, Inches(0.4),
                 role, size=11, italic=True, bold=True, color=color,
                 align=PP_ALIGN.CENTER)
        # divider
        add_rect(slide, x + Inches(0.5), y0 + Inches(2.7), cw - Inches(1.0), Inches(0.02),
                 fill=color)
        # contributions
        runs = []
        for j, item in enumerate(items):
            if j > 0:
                runs.append({"newline": True})
            runs.append({"text": "● ", "size": 11, "color": color, "bold": True})
            runs.append({"text": item, "size": 10, "color": DARK})
        add_multi_text(slide, x + Inches(0.25), y0 + Inches(2.85),
                       cw - Inches(0.5), Inches(2.0), runs)


# =========================================================================
# SLIDE 4 — PROJECT OVERVIEW
# =========================================================================
def slide_overview():
    slide = slide_chrome("Project Overview", "From raw crowdsourced salaries to an interactive Tableau master dashboard", page_n=4)

    # Left: big paragraph
    add_text(slide, Inches(0.6), Inches(1.6), Inches(7.0), Inches(0.5),
             "What we built", size=18, bold=True, color=FOREST)
    add_text(slide, Inches(0.6), Inches(2.1), Inches(7.0), Inches(3.0),
             "An end-to-end data visualization project that converts a real, "
             "first-party salary dataset (147,348 records · 87 countries · "
             "2022–2025) into an interactive Tableau dashboard. The work "
             "spans a Python preprocessing pipeline, twenty exploratory "
             "figures, ten coordinated Tableau worksheets, and a 40-page "
             "report. The deliverable lets a stakeholder click a country on "
             "the map and watch every other panel re-scope.",
             size=12, color=DARK, line_spacing=1.4)

    # Objective bullets
    add_text(slide, Inches(0.6), Inches(4.7), Inches(7.0), Inches(0.4),
             "Objectives", size=18, bold=True, color=FOREST)
    objectives = [
        ("◆", "Data Integrity",      "Clean and normalise crowdsourced salaries"),
        ("◆", "Feature Engineering", "Add 12 derived columns for richer encodings"),
        ("◆", "Exploratory Analysis","20 matplotlib/seaborn figures profiling every facet"),
        ("◆", "Dashboarding",        "10 worksheets + 5 KPIs + 4 insights, 1200×1500 canvas"),
    ]
    y = Inches(5.2)
    for icon, label, desc in objectives:
        add_multi_text(slide, Inches(0.6), y, Inches(7.0), Inches(0.45), runs=[
            {"text": icon + "  ", "size": 13, "bold": True, "color": EMERALD},
            {"text": label + " — ", "size": 12, "bold": True, "color": FOREST},
            {"text": desc, "size": 11, "color": DARK},
        ])
        y += Inches(0.4)

    # Right: KPI tiles
    kpi_x = Inches(8.2); kpi_w = Inches(2.4); kpi_h = Inches(1.05); gap = Inches(0.15)
    tiles = [
        ("147,348",  "TOTAL RECORDS",  FOREST),
        ("$147,000", "MEDIAN SALARY",  EMERALD),
        ("87",       "COUNTRIES",      JADE),
        ("10",       "WORKSHEETS",     SAGE),
        ("4",        "YEARS COVERED",  HONEY),
        ("1.91×",    "EXEC vs ENTRY",  RUST),
    ]
    cols = 2; rows = 3
    for i, (val, lbl, c) in enumerate(tiles):
        col = i % cols; row = i // cols
        x = kpi_x + col * (kpi_w + gap)
        y = Inches(1.6) + row * (kpi_h + gap)
        kpi_tile(slide, x, y, kpi_w, kpi_h, val, lbl, c)


# =========================================================================
# SLIDE 5 — 5 GUIDING QUESTIONS
# =========================================================================
def slide_questions():
    slide = slide_chrome("Five Guiding Questions", "What the dashboard is engineered to answer", page_n=5)

    questions = [
        ("01", "Which AI/ML/Data Science roles pay the most?",
               "ML Eng. / MLOps leads at $194K — 81% above Data Analyst.", FOREST),
        ("02", "How does salary grow from Entry Level to Executive?",
               "Entry $103K → Mid $140K → Senior $170K → Exec $196K  (1.91×)", EMERALD),
        ("03", "How has remote work changed year over year?",
               "Remote dropped from 53.5% (2022) to 20.3% (2025).", JADE),
        ("04", "Which countries lead global AI hiring?",
               "US captures 90.4% of submissions; North America $159K avg.", HONEY),
        ("05", "What does the salary distribution look like?",
               "$100K–$200K holds 55% of all records; $260K+ holds 7.3%.", RUST),
    ]
    y0 = Inches(1.7); h = Inches(1.0); gap = Inches(0.1)
    for i, (num, q, a, color) in enumerate(questions):
        y = y0 + i * (h + gap)
        # left vertical accent
        add_rect(slide, Inches(0.45), y, Inches(0.15), h, fill=color)
        card(slide, Inches(0.65), y, SW - Inches(1.1), h, fill=WHITE, border=color)
        add_text(slide, Inches(0.85), y + Inches(0.05), Inches(1.0), h - Inches(0.1),
                 num, size=28, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(2.0), y + Inches(0.10), Inches(10), Inches(0.45),
                 q, size=15, bold=True, color=FOREST)
        add_text(slide, Inches(2.0), y + Inches(0.50), Inches(10), Inches(0.45),
                 "→  " + a, size=11, italic=True, color=GRAY)


# =========================================================================
# SLIDE 6 — ABOUT THE DATASET
# =========================================================================
def slide_dataset():
    slide = slide_chrome("About the Dataset", "Real, crowdsourced AI/ML/Data Science salary survey  ·  CC0 Public Domain", page_n=6)

    # Left description
    add_text(slide, Inches(0.6), Inches(1.6), Inches(7.5), Inches(0.4),
             "What it is", size=18, bold=True, color=FOREST)
    add_text(slide, Inches(0.6), Inches(2.05), Inches(7.5), Inches(2.5),
             "Global Salaries in AI, ML, Data Science and Big Data — a "
             "crowdsourced compensation survey published by aijobs.net "
             "(the same first-party organisation that runs the live job "
             "board). Every row is a single anonymous salary submission "
             "covering 2022–2025. We standardise 406 free-text job "
             "titles into 15 categories and add 12 engineered analytical "
             "columns, ending at 147,348 rows × 27 columns.",
             size=12, color=DARK, line_spacing=1.4)

    # 4 column boxes for raw fields
    add_text(slide, Inches(0.6), Inches(4.6), Inches(7.5), Inches(0.4),
             "What each row carries", size=18, bold=True, color=FOREST)
    raw_boxes = [
        ("WHEN", ["work_year"], JADE),
        ("WHO",  ["experience_level", "exp_label", "employee_residence"], EMERALD),
        ("WHAT", ["job_title", "title_group", "company_size"], FOREST),
        ("HOW",  ["salary_in_usd", "remote_ratio", "work_mode"], HONEY),
    ]
    bx = Inches(0.6); by = Inches(5.05)
    bw = Inches(1.85); bh = Inches(1.6); gap = Inches(0.1)
    for i, (head, fields, color) in enumerate(raw_boxes):
        x = bx + i * (bw + gap)
        card(slide, x, by, bw, bh, fill=WHITE, border=color)
        add_rect(slide, x, by, bw, Inches(0.35), fill=color)
        add_text(slide, x, by, bw, Inches(0.35),
                 head, size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        for j, fld in enumerate(fields):
            add_text(slide, x + Inches(0.1), by + Inches(0.45) + j*Inches(0.3),
                     bw - Inches(0.2), Inches(0.3),
                     "● " + fld, size=10, color=DARK)

    # Right side: properties table
    add_text(slide, Inches(8.5), Inches(1.6), Inches(4.5), Inches(0.4),
             "Headline properties", size=18, bold=True, color=FOREST)

    props = [
        ("Records (cleaned)",           "147,348"),
        ("Original raw rows",           "151,445"),
        ("Time period",                 "2022 – 2025"),
        ("Countries",                   "87"),
        ("Role categories",             "15  (from 406 raw)"),
        ("Salary range",                "$37,974 – $385,000"),
        ("Median / Mean salary",        "$147,000 / $156,237"),
        ("Skewness",                    "0.68 (right-skewed)"),
        ("Final columns",               "27 (11 raw + 4 label + 12 engineered)"),
        ("License",                     "CC0 1.0 Public Domain"),
    ]
    table_y = Inches(2.1)
    row_h = Inches(0.42)
    for k, v in props:
        add_rect(slide, Inches(8.5), table_y, Inches(4.5), row_h, fill=LIGHT_MINT)
        add_text(slide, Inches(8.6), table_y, Inches(2.6), row_h,
                 k, size=10, bold=True, color=FOREST,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(11.0), table_y, Inches(1.95), row_h,
                 v, size=10, color=DARK, align=PP_ALIGN.RIGHT,
                 anchor=MSO_ANCHOR.MIDDLE)
        table_y += row_h + Inches(0.04)


# =========================================================================
# SLIDE 7 — DATASET AUTHENTICITY
# =========================================================================
def slide_authenticity():
    slide = slide_chrome("Why This Dataset Is Real", "Authenticity markers we verified before any preprocessing", page_n=7)

    # Verification screenshots (4 thumbnails)
    pics = [
        ("pic1_foorilla_hiring.png", "Foorilla hiring platform"),
        ("pic2_foorilla_media.png",  "Media + industry recognition"),
        ("pic3_foorilla_insight.png","Internal market insight tools"),
        ("pic4_aijobs_search.png",   "aijobs.net live search UI"),
    ]
    pw = Inches(2.85); ph = Inches(1.8); gap = Inches(0.15)
    px0 = Inches(0.5); py = Inches(1.6)
    for i, (rel, lbl) in enumerate(pics):
        x = px0 + i * (pw + gap)
        path = ASSETS / rel
        if path.exists():
            slide.shapes.add_picture(str(path), x, py, width=pw, height=ph)
        else:
            add_rect(slide, x, py, pw, ph, fill=LIGHT_MINT, line=EMERALD, line_w=2)
        add_text(slide, x, py + ph + Inches(0.05), pw, Inches(0.3),
                 lbl, size=10, italic=True, color=GRAY, align=PP_ALIGN.CENTER)

    # Bullet markers below
    bullets = [
        ("●", FOREST,  "First-party Kaggle release — published by the official aijobs organisation account"),
        ("●", EMERALD, "Real-world distribution: US 90.4% of submissions matches public hiring statistics"),
        ("●", JADE,    "Right-skewed salary distribution (skewness 0.68) — characteristic of real labour markets"),
        ("●", SAGE,    "406 distinct raw job-title strings reflect genuine free-text variation, not synthetic data"),
        ("●", HONEY,   "Updated weekly by the platform that runs the live aijobs.net job board"),
        ("●", RUST,    "99.4% full-time employment ratio matches industry averages"),
    ]
    by = Inches(4.0); h = Inches(0.4)
    for i, (icon, c, text) in enumerate(bullets):
        col = i % 2; row = i // 2
        x = Inches(0.5) + col * Inches(6.3)
        y = by + row * Inches(0.55)
        add_multi_text(slide, x, y, Inches(6.2), h, runs=[
            {"text": icon + "  ", "size": 14, "bold": True, "color": c},
            {"text": text, "size": 11, "color": DARK},
        ], anchor=MSO_ANCHOR.MIDDLE)

    # Bottom callout
    by = Inches(6.2)
    card(slide, Inches(0.5), by, SW - Inches(1.0), Inches(0.6),
         fill=PEACH, border=HONEY)
    add_multi_text(slide, Inches(0.7), by, Inches(12), Inches(0.6), runs=[
        {"text": "✦  Conclusion:  ", "size": 13, "bold": True, "color": HONEY},
        {"text": "every row in our final ", "size": 12, "color": DARK},
        {"text": "salaries_enhanced.csv ", "size": 12, "bold": True, "color": FOREST, "font": "Consolas"},
        {"text": "is a real salary submission. No synthetic data is introduced at any stage of the pipeline.",
         "size": 12, "color": DARK},
    ], anchor=MSO_ANCHOR.MIDDLE)


# =========================================================================
# SLIDE 8 — PREPROCESSING PIPELINE OVERVIEW
# =========================================================================
def slide_pipeline_overview():
    slide = slide_chrome("Preprocessing Pipeline", "Five Python scripts: 151,445 rows in → 147,348 × 27 cols out", page_n=8)

    stages = [
        ("01_explore.py",            "EXPLORE",         "Profile shape · nulls · skew",  FOREST),
        ("02_clean.py",              "CLEAN",           "FT-only · 2022-25 · 1-99 pct",  EMERALD),
        ("03_title_grouping.py",     "STANDARDIZE",     "406 raw titles → 15 groups",    JADE),
        ("04_summary_stats.py",      "SUMMARIZE",       "10 KPI / category CSVs",        SAGE),
        ("05_feature_engineering.py","ENGINEER",        "12 derived analytical cols",    HONEY),
    ]
    sw = Inches(2.4); sh_ = Inches(2.6); gap = Inches(0.15); x0 = Inches(0.6); y0 = Inches(1.8)

    for i, (script, head, desc, color) in enumerate(stages):
        x = x0 + i * (sw + gap)
        # arrow connector between stages
        if i > 0:
            ax = x - gap; ay = y0 + sh_/2 - Inches(0.1)
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                           ax - Inches(0.05), ay,
                                           gap + Inches(0.1), Inches(0.2))
            arrow.shadow.inherit = False
            arrow.fill.solid(); arrow.fill.fore_color.rgb = color
            arrow.line.fill.background()
        card(slide, x, y0, sw, sh_, fill=WHITE, border=color)
        # Header
        add_rect(slide, x, y0, sw, Inches(0.5), fill=color)
        add_text(slide, x, y0, sw, Inches(0.5),
                 f"{i+1:02d}  ·  {head}",
                 size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Script name
        add_text(slide, x, y0 + Inches(0.65), sw, Inches(0.4),
                 script, size=11, bold=True, color=FOREST,
                 align=PP_ALIGN.CENTER, font="Consolas")
        # description
        add_text(slide, x + Inches(0.15), y0 + Inches(1.15), sw - Inches(0.3), Inches(0.7),
                 desc, size=11, color=DARK, align=PP_ALIGN.CENTER, line_spacing=1.4,
                 anchor=MSO_ANCHOR.MIDDLE)
        # footer with status
        add_rect(slide, x + Inches(0.2), y0 + sh_ - Inches(0.4),
                 sw - Inches(0.4), Inches(0.25), fill=LIGHT_MINT)
        add_text(slide, x, y0 + sh_ - Inches(0.4), sw, Inches(0.25),
                 "✓ DONE", size=9, bold=True, color=EMERALD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Bottom row counts
    by = Inches(4.7)
    counts = [
        ("RAW",    "151,445", FOREST),
        ("FT only","150,541", EMERALD),
        ("2022-25","150,264", JADE),
        ("Outlier-free","147,348", SAGE),
        ("FINAL",  "147,348 × 27 cols", HONEY),
    ]
    cw_ = Inches(2.4); ch_ = Inches(1.3); gap_ = Inches(0.15)
    for i, (lbl, val, c) in enumerate(counts):
        x = Inches(0.6) + i * (cw_ + gap_)
        card(slide, x, by, cw_, ch_, fill=CREAM, border=c)
        add_text(slide, x, by + Inches(0.15), cw_, Inches(0.45),
                 val, size=18, bold=True, color=c,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x, by + Inches(0.7), cw_, Inches(0.4),
                 lbl, size=11, color=GRAY, align=PP_ALIGN.CENTER)

    # Caption
    add_text(slide, Inches(0.6), Inches(6.15), SW - Inches(1.2), Inches(0.5),
             "We removed 4,097 rows total (2.7%): 904 non-full-time + 277 outside 2022–2025 + "
             "2,916 outliers below the 1st or above the 99th percentile.",
             size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# =========================================================================
# SLIDE 9 — STAGE 1 + 2 (EXPLORE + CLEAN)  — WHY
# =========================================================================
def slide_stage_clean():
    slide = slide_chrome("Stages 1 & 2 — Explore  ·  Clean",
                         "Why the cleaning rules exist", page_n=9)

    # Left: explore
    add_rect(slide, Inches(0.5), Inches(1.6), Inches(0.15), Inches(4.0), fill=FOREST)
    add_text(slide, Inches(0.75), Inches(1.6), Inches(6.0), Inches(0.5),
             "01  ·  EXPLORE", size=18, bold=True, color=FOREST)
    add_text(slide, Inches(0.75), Inches(2.1), Inches(6.0), Inches(0.4),
             "Purpose", size=12, bold=True, color=EMERALD)
    add_text(slide, Inches(0.75), Inches(2.4), Inches(6.0), Inches(2.0),
             "Profile the raw 151,445-row CSV before touching it. Confirm "
             "shape, datatypes, missing values, salary distribution, "
             "skewness, and the spread of unique values per column.",
             size=11, color=DARK, line_spacing=1.4)
    add_text(slide, Inches(0.75), Inches(3.5), Inches(6.0), Inches(0.4),
             "Why it matters", size=12, bold=True, color=EMERALD)
    add_text(slide, Inches(0.75), Inches(3.85), Inches(6.0), Inches(2.0),
             "Surfaced the right-skew (0.68), 0 missing values, 99.4% "
             "full-time, and 406 unique title strings — the evidence base "
             "for every later cleaning decision.",
             size=11, color=DARK, line_spacing=1.4)

    # Right: cleaning
    add_rect(slide, Inches(7.0), Inches(1.6), Inches(0.15), Inches(4.0), fill=EMERALD)
    add_text(slide, Inches(7.25), Inches(1.6), Inches(6.0), Inches(0.5),
             "02  ·  CLEAN", size=18, bold=True, color=EMERALD)

    # Three cleaning steps as mini-cards
    steps = [
        ("Keep Full-Time",         "151,445 → 150,541",  "−904",   "Removes contractor / part-time noise so the salary scale is comparable."),
        ("Keep 2022 – 2025",       "150,541 → 150,264",  "−277",   "Pre-2022 sample is too small (~1k records) to draw stable trend lines."),
        ("Outlier removal (1–99%)","150,264 → 147,348",  "−2,916", "Bounds: $37,974 – $385,000. Trims data-entry errors and extreme top-of-market."),
    ]
    sy = Inches(2.1)
    for i, (title, count, removed, why) in enumerate(steps):
        y = sy + i * Inches(1.15)
        card(slide, Inches(7.25), y, Inches(5.7), Inches(1.05), fill=WHITE, border=EMERALD)
        add_text(slide, Inches(7.45), y + Inches(0.05), Inches(3.0), Inches(0.4),
                 title, size=12, bold=True, color=FOREST)
        add_text(slide, Inches(7.45), y + Inches(0.4), Inches(2.5), Inches(0.3),
                 count, size=10, color=DARK, font="Consolas")
        add_text(slide, Inches(10.5), y + Inches(0.05), Inches(2.4), Inches(0.4),
                 removed, size=14, bold=True, color=RUST,
                 align=PP_ALIGN.RIGHT)
        add_text(slide, Inches(7.45), y + Inches(0.65), Inches(5.4), Inches(0.4),
                 why, size=10, italic=True, color=GRAY)

    # Bottom: bounds callout
    by = Inches(6.0)
    card(slide, Inches(0.5), by, SW - Inches(1.0), Inches(0.7),
         fill=LIGHT_MINT, border=SAGE)
    add_multi_text(slide, Inches(0.7), by, Inches(12), Inches(0.7), runs=[
        {"text": "Final salary range after cleaning:  ", "size": 13, "bold": True, "color": EMERALD},
        {"text": "$37,974 ", "size": 13, "bold": True, "color": FOREST},
        {"text": "to ", "size": 13, "color": DARK},
        {"text": "$385,000 ", "size": 13, "bold": True, "color": FOREST},
        {"text": " ·  4,097 rows removed total (2.7% of raw)", "size": 13, "color": GRAY},
    ], anchor=MSO_ANCHOR.MIDDLE)


# =========================================================================
# SLIDE 10 — STAGE 3 TITLE GROUPING
# =========================================================================
def slide_stage_grouping():
    slide = slide_chrome("Stage 3 — Title Grouping", "406 raw free-text titles → 15 standardized categories", page_n=10)

    # Why & technique
    add_text(slide, Inches(0.5), Inches(1.6), Inches(6.5), Inches(0.5),
             "Technique  ·  keyword string matching", size=14, bold=True, color=EMERALD)
    add_text(slide, Inches(0.5), Inches(2.1), Inches(6.5), Inches(2.0),
             "Each free-text title is lowercased and searched for "
             "category-specific keyword sets. \"ML Eng.\", \"Machine "
             "Learning Engineer\", and \"MLOps Engineer\" all collapse "
             "into one canonical category: ML Engineer / MLOps.",
             size=11, color=DARK, line_spacing=1.4)
    add_text(slide, Inches(0.5), Inches(3.5), Inches(6.5), Inches(0.4),
             "Why it matters", size=14, bold=True, color=EMERALD)
    add_text(slide, Inches(0.5), Inches(3.95), Inches(6.5), Inches(2.5),
             "Without grouping, the role bar chart would be unreadable "
             "(406 categories) and statistical comparisons would be "
             "diluted. The 15-bucket schema preserves real signal — ML "
             "Engineer ($194K) is genuinely distinct from Data Analyst "
             "($107K) — while keeping every chart legible.",
             size=11, color=DARK, line_spacing=1.4)

    # Right: example mapping
    add_text(slide, Inches(7.5), Inches(1.6), Inches(5.3), Inches(0.5),
             "Example mappings", size=14, bold=True, color=EMERALD)
    examples = [
        ("\"machine learning engineer\"",  "→",  "ML Engineer / MLOps"),
        ("\"mlops engineer\"",             "→",  "ML Engineer / MLOps"),
        ("\"applied scientist\"",          "→",  "Data Scientist"),
        ("\"data engineering manager\"",   "→",  "Data Engineer"),
        ("\"bi analyst\"",                 "→",  "Data Analyst"),
        ("\"deep learning researcher\"",   "→",  "AI / ML Researcher"),
        ("\"head of analytics\"",          "→",  "Manager / Director"),
    ]
    ey = Inches(2.1); rh = Inches(0.45)
    for raw, arrow, cat in examples:
        add_rect(slide, Inches(7.5), ey, Inches(5.3), rh, fill=LIGHT_MINT)
        add_text(slide, Inches(7.6), ey, Inches(2.6), rh,
                 raw, size=10, color=GRAY, anchor=MSO_ANCHOR.MIDDLE,
                 font="Consolas", italic=True)
        add_text(slide, Inches(10.1), ey, Inches(0.3), rh,
                 arrow, size=12, bold=True, color=HONEY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(10.4), ey, Inches(2.3), rh,
                 cat, size=10, bold=True, color=FOREST,
                 anchor=MSO_ANCHOR.MIDDLE)
        ey += rh + Inches(0.04)

    # Bottom band: top role pay
    by = Inches(6.4)
    card(slide, Inches(0.5), by, SW - Inches(1.0), Inches(0.6), fill=PEACH, border=HONEY)
    add_multi_text(slide, Inches(0.7), by, Inches(12), Inches(0.6), runs=[
        {"text": "Outcome:  ", "size": 13, "bold": True, "color": HONEY},
        {"text": "ML Engineer / MLOps  $194,303", "size": 12, "bold": True, "color": FOREST},
        {"text": "  →  ", "size": 12, "color": GRAY},
        {"text": "Software Engineer  $180,027", "size": 12, "color": DARK},
        {"text": "  →  ", "size": 12, "color": GRAY},
        {"text": "Data Analyst  $107,465", "size": 12, "color": DARK},
    ], anchor=MSO_ANCHOR.MIDDLE)


# =========================================================================
# SLIDE 11 — STAGES 4 & 5
# =========================================================================
def slide_stages_4_5():
    slide = slide_chrome("Stages 4 & 5 — Summarize  ·  Engineer 12 columns",
                         "From flat data to analysis-ready features for Tableau", page_n=11)

    # left summary stats
    add_rect(slide, Inches(0.5), Inches(1.6), Inches(0.15), Inches(4.5), fill=SAGE)
    add_text(slide, Inches(0.75), Inches(1.6), Inches(6), Inches(0.5),
             "04  ·  SUMMARIZE", size=18, bold=True, color=SAGE)
    add_text(slide, Inches(0.75), Inches(2.1), Inches(6), Inches(0.4),
             "Output", size=12, bold=True, color=EMERALD)
    add_text(slide, Inches(0.75), Inches(2.45), Inches(6), Inches(0.3),
             "10 summary CSVs in notebooks/outputs/",
             size=11, color=DARK, font="Consolas")
    items = [
        "summary_year_trend.csv",
        "summary_by_role.csv",
        "summary_by_experience.csv",
        "summary_by_country.csv",
        "summary_by_region.csv",
        "summary_us_vs_row.csv",
        "summary_ai_core_vs_tech.csv",
        "summary_remote_trend.csv",
        "summary_salary_bands.csv",
        "analysis_summary.csv",
    ]
    iy = Inches(2.85)
    for it in items:
        add_text(slide, Inches(0.95), iy, Inches(5.5), Inches(0.25),
                 "● " + it, size=10, color=DARK, font="Consolas")
        iy += Inches(0.27)

    # right engineering
    add_rect(slide, Inches(7.0), Inches(1.6), Inches(0.15), Inches(4.5), fill=HONEY)
    add_text(slide, Inches(7.25), Inches(1.6), Inches(6), Inches(0.5),
             "05  ·  ENGINEER (12 derived columns)",
             size=18, bold=True, color=HONEY)

    eng_groups = [
        ("◆ CATEGORICAL BINS", "salary_band  ·  salary_quartile",            FOREST),
        ("◆ GEOGRAPHIC",        "region  ·  is_us",                            EMERALD),
        ("◆ ROLE FAMILY",       "is_ai_core",                                  JADE),
        ("◆ SORT KEY",          "exp_order",                                   SAGE),
        ("◆ YEAR-RELATIVE",     "year_avg_salary  ·  salary_vs_year_avg  ·  …pct", HONEY),
        ("◆ ROLE-RELATIVE",     "role_avg_salary  ·  salary_vs_role_avg  ·  …pct", RUST),
    ]
    ey = Inches(2.15)
    for head, fields, c in eng_groups:
        card(slide, Inches(7.25), ey, Inches(5.7), Inches(0.65), fill=WHITE, border=c)
        add_rect(slide, Inches(7.25), ey, Inches(0.15), Inches(0.65), fill=c)
        add_text(slide, Inches(7.5), ey + Inches(0.05), Inches(5.5), Inches(0.3),
                 head, size=11, bold=True, color=c)
        add_text(slide, Inches(7.5), ey + Inches(0.32), Inches(5.5), Inches(0.3),
                 fields, size=10, color=DARK, font="Consolas")
        ey += Inches(0.72)


# =========================================================================
# SLIDE 12 — EDA SAMPLE
# =========================================================================
def slide_eda_overview():
    slide = slide_chrome("Exploratory Data Analysis", "20 figures across 7 themes — built in Python before designing Tableau", page_n=12)

    pics = [
        ("01_salary_distribution.png",      "Salary distribution",  EMERALD),
        ("04_salary_by_role.png",            "Top-paying roles",     FOREST),
        ("10_remote_trend_stacked.png",      "Remote-work collapse", RUST),
        ("11_salary_by_region.png",          "Regional pay gap",     HONEY),
        ("07_salary_vs_experience.png",      "Career progression",   JADE),
        ("13_ai_core_vs_data_tech.png",      "AI Core premium",      SAGE),
    ]
    cols = 3; rows = 2
    pw = Inches(4.0); ph = Inches(2.4); gap_x = Inches(0.15); gap_y = Inches(0.25)
    x0 = Inches(0.5); y0 = Inches(1.6)
    for i, (fname, lbl, c) in enumerate(pics):
        col = i % cols; row = i // cols
        x = x0 + col * (pw + gap_x); y = y0 + row * (ph + gap_y + Inches(0.3))
        card(slide, x, y, pw, ph + Inches(0.3), fill=WHITE, border=c)
        # Image inside
        path = (FIG / fname)
        if path.exists():
            slide.shapes.add_picture(str(path),
                                     x + Inches(0.05), y + Inches(0.05),
                                     width=pw - Inches(0.1), height=ph - Inches(0.05))
        # Label strip
        add_rect(slide, x, y + ph, pw, Inches(0.3), fill=c)
        add_text(slide, x, y + ph, pw, Inches(0.3),
                 lbl, size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# =========================================================================
# SLIDE 13 — TABLEAU MASTER DASHBOARD
# =========================================================================
def slide_tableau_master():
    slide = slide_chrome("Tableau Master Dashboard", "1200 × 1500  ·  10 worksheets  ·  5 KPI tiles  ·  4 insight callouts", page_n=13)

    pic = FIG / "Tableau_Master_Dashboard.png"
    if pic.exists():
        # Full image, centered
        slide.shapes.add_picture(str(pic),
                                 Inches(2.0), Inches(1.55),
                                 width=Inches(9.0), height=Inches(5.55))

    # Floating side annotations
    sy = Inches(1.7)
    notes = [
        ("✦ Forest emerald palette",          EMERALD),
        ("✦ Click-to-filter actions",          JADE),
        ("✦ 5 colour-coded KPI tiles",         HONEY),
        ("✦ 4 insight callouts",               RUST),
    ]
    for txt, c in notes:
        add_text(slide, Inches(0.4), sy, Inches(1.6), Inches(0.4),
                 txt, size=11, bold=True, color=c)
        sy += Inches(0.5)

    sy = Inches(1.7)
    notes2 = [
        ("Top-down narrative grid →",          EMERALD),
        ("Section bands group sheets →",       JADE),
        ("Mint section headers →",             HONEY),
        ("Forest title + footer bars →",       RUST),
    ]
    for txt, c in notes2:
        add_text(slide, Inches(11.2), sy, Inches(2.0), Inches(0.4),
                 txt, size=11, bold=True, color=c, align=PP_ALIGN.RIGHT)
        sy += Inches(0.5)


# =========================================================================
# SLIDES 14 / 15 — SHEET GROUPS
# =========================================================================
def slide_sheets_group_1():
    slide = slide_chrome("Section 1 — Salary Trends & Role Compensation",
                         "Sheet 1 — Salary Trend  ·  Sheet 2 — Salary By Role", page_n=14)

    # Two charts side by side
    pics = [
        ("Tableau_Salary_Trend.png",   "Salary Trend (2022 – 2025)",
         "Median climbed 46% from $102K (2022) → $149K (2024). Submission volume grew 48× as the platform matured."),
        ("Tableau_Salary_By_Role.png", "Salary By Role (Top 10)",
         "ML Engineer / MLOps leads at $194K — 81% above Data Analyst. Top six roles all clear the $170K mark."),
    ]
    cw = Inches(6.0); ch = Inches(3.0); gap = Inches(0.3); x0 = Inches(0.5); y0 = Inches(1.7)
    for i, (fname, title, analysis) in enumerate(pics):
        x = x0 + i * (cw + gap)
        # Title
        add_text(slide, x, y0, cw, Inches(0.4),
                 title, size=14, bold=True, color=FOREST)
        # Image
        path = FIG / fname
        if path.exists():
            slide.shapes.add_picture(str(path),
                                     x, y0 + Inches(0.45),
                                     width=cw, height=ch)
        # Analysis card
        ay = y0 + Inches(0.45) + ch + Inches(0.2)
        card(slide, x, ay, cw, Inches(1.4), fill=PEACH, border=HONEY)
        add_multi_text(slide, x + Inches(0.15), ay + Inches(0.1),
                       cw - Inches(0.3), Inches(1.2), runs=[
            {"text": "ANALYSIS  ·  ", "size": 11, "bold": True, "color": HONEY},
            {"text": analysis, "size": 11, "color": DARK},
        ])


def slide_sheets_group_2():
    slide = slide_chrome("Section 2 — Work Mode  ·  Experience  ·  Geography",
                         "Sheets 3, 4, 5 — Work Mode donut, Experience bars, World choropleth", page_n=15)

    # 3 charts in row
    pics = [
        ("Tableau_Work_Mode.png",          "Work Mode (Donut)",
         "On-site dominates at ~74%. Remote shrunk; hybrid is gone."),
        ("Tableau_Salary_Vs_Experience.png", "Salary vs Experience",
         "$103K → $140K → $170K → $196K  (1.91× from Entry to Exec)."),
        ("Tableau_Global_Map.png",          "Global Hiring Map",
         "US dominates volume. Click any country to filter every other panel."),
    ]
    cw = Inches(4.0); ch = Inches(2.8); gap = Inches(0.15); x0 = Inches(0.5); y0 = Inches(1.7)
    for i, (fname, title, txt) in enumerate(pics):
        x = x0 + i * (cw + gap)
        add_text(slide, x, y0, cw, Inches(0.4),
                 title, size=13, bold=True, color=FOREST)
        path = FIG / fname
        if path.exists():
            slide.shapes.add_picture(str(path),
                                     x, y0 + Inches(0.45),
                                     width=cw, height=ch)
        ay = y0 + Inches(0.45) + ch + Inches(0.15)
        card(slide, x, ay, cw, Inches(1.6), fill=PEACH, border=HONEY)
        add_multi_text(slide, x + Inches(0.15), ay + Inches(0.1),
                       cw - Inches(0.3), Inches(1.4), runs=[
            {"text": "ANALYSIS  ·  ", "size": 10, "bold": True, "color": HONEY},
            {"text": txt, "size": 10, "color": DARK},
        ])


def slide_sheets_group_3():
    slide = slide_chrome("Section 3 — Regional · Company Size · Role Family · Market Concentration",
                         "Sheets 6, 7, 8, 9 — four comparative panels added in the latest rev", page_n=16)

    pics = [
        ("Tableau_Salary_By_Region.png",       "By Region",
         "North America $159K dominates. 94% premium over Latin America."),
        ("Tableau_Salary_By_Company_Size.png", "By Company Size",
         "Small $103K → Medium $155K → Large $165K. The big jump is Small→Medium."),
        ("Tableau_AI_Core_Vs_Data_Tech.png",   "AI Core vs Data&Tech",
         "AI/ML Core $171K vs Data&Tech $151K — a 13% \"AI tax\" premium."),
        ("Tableau_US_Vs_Rest_Of_World.png",    "US vs Rest of World",
         "US $161K vs RoW $108K — 49% premium per record (133K vs 14K rows)."),
    ]
    cw = Inches(2.95); ch = Inches(2.4); gap = Inches(0.15); x0 = Inches(0.5); y0 = Inches(1.7)
    for i, (fname, title, txt) in enumerate(pics):
        x = x0 + i * (cw + gap)
        add_text(slide, x, y0, cw, Inches(0.4),
                 title, size=12, bold=True, color=FOREST)
        path = FIG / fname
        if path.exists():
            slide.shapes.add_picture(str(path),
                                     x, y0 + Inches(0.45),
                                     width=cw, height=ch)
        ay = y0 + Inches(0.45) + ch + Inches(0.15)
        card(slide, x, ay, cw, Inches(1.7), fill=PEACH, border=HONEY)
        add_multi_text(slide, x + Inches(0.1), ay + Inches(0.1),
                       cw - Inches(0.2), Inches(1.5), runs=[
            {"text": "● ", "size": 10, "bold": True, "color": HONEY},
            {"text": txt, "size": 10, "color": DARK},
        ])


def slide_sheets_group_4():
    slide = slide_chrome("Section 4 — Salary Distribution Across Pay Bands",
                         "Sheet 10 — full-width histogram of every record across 6 bands", page_n=17)

    path = FIG / "Tableau_Salary_Band_Distribution.png"
    if path.exists():
        slide.shapes.add_picture(str(path),
                                 Inches(0.7), Inches(1.7),
                                 width=Inches(11.9), height=Inches(2.5))

    # Six labelled stat tiles below
    bands = [
        ("<$60K",      "3.9%",   "5,742",  RUST),
        ("$60K–$100K", "18.0%",  "26,572", HONEY),
        ("$100K–$150K","31.0%",  "45,617", FOREST),
        ("$150K–$200K","24.3%",  "35,738", JADE),
        ("$200K–$260K","15.5%",  "22,901", EMERALD),
        ("$260K+",     "7.3%",   "10,778", SAGE),
    ]
    bx = Inches(0.5); by = Inches(4.5)
    bw = Inches(2.05); bh = Inches(1.3); gap = Inches(0.1)
    for i, (lbl, pct, n, c) in enumerate(bands):
        x = bx + i * (bw + gap)
        card(slide, x, by, bw, bh, fill=WHITE, border=c)
        add_rect(slide, x, by, bw, Inches(0.3), fill=c)
        add_text(slide, x, by, bw, Inches(0.3),
                 lbl, size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x, by + Inches(0.35), bw, Inches(0.5),
                 pct, size=22, bold=True, color=c,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x, by + Inches(0.9), bw, Inches(0.35),
                 n + " rows", size=9, color=GRAY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Insight callout
    iy = Inches(6.05)
    card(slide, Inches(0.5), iy, SW - Inches(1.0), Inches(0.7),
         fill=PEACH, border=HONEY)
    add_multi_text(slide, Inches(0.7), iy, Inches(12), Inches(0.7), runs=[
        {"text": "✦  Insight  ·  ", "size": 13, "bold": True, "color": HONEY},
        {"text": "55% of all records sit in the $100K–$200K window.", "size": 12, "bold": True, "color": FOREST},
        {"text": "  Modal class is $100K–$150K with 31% / ~46K rows. The top tier $260K+ holds a non-trivial 7.3% / 10,778 rows.",
         "size": 12, "color": DARK},
    ], anchor=MSO_ANCHOR.MIDDLE)


# =========================================================================
# SLIDE 18 — STORY THE DASHBOARD TELLS
# =========================================================================
def slide_story():
    slide = slide_chrome("The Story Our Dashboard Tells",
                         "A four-act narrative running top-to-bottom on the canvas", page_n=18)

    acts = [
        ("ACT 1", "The market is broadening, not inflating",
         "Volume jumps 48× while median salary plateaus near $146K.",
         FOREST,  "Sheets 1, 2"),
        ("ACT 2", "Career level matters more than work mode",
         "$103K → $196K from Entry to Exec, while remote retreats to 20%.",
         EMERALD, "Sheets 3, 5"),
        ("ACT 3", "Geography concentrates the money in North America",
         "US 90% of submissions at $161K; world average drops to $108K.",
         JADE,    "Sheets 4, 6, 9"),
        ("ACT 4", "Direct AI/ML exposure is rewarded",
         "AI Core premium of 13%; ML Engineer / MLOps tops the pay table.",
         HONEY,   "Sheets 7, 8, 10"),
    ]

    y0 = Inches(1.7); h = Inches(1.25); gap = Inches(0.1)
    for i, (act, headline, body, color, ref) in enumerate(acts):
        y = y0 + i * (h + gap)
        # Vertical accent
        add_rect(slide, Inches(0.5), y, Inches(0.18), h, fill=color)
        # Card
        card(slide, Inches(0.7), y, SW - Inches(1.2), h, fill=WHITE, border=color)
        # Act label
        add_text(slide, Inches(0.85), y + Inches(0.05), Inches(1.0), Inches(0.5),
                 act, size=22, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)
        # Headline
        add_text(slide, Inches(2.0), y + Inches(0.1), Inches(8.5), Inches(0.45),
                 headline, size=15, bold=True, color=FOREST)
        # Body
        add_text(slide, Inches(2.0), y + Inches(0.55), Inches(8.5), Inches(0.5),
                 body, size=11, color=DARK)
        # Refs on the right
        add_text(slide, Inches(10.5), y + Inches(0.1), Inches(2.0), Inches(0.4),
                 "REFERS TO", size=8, bold=True, color=GRAY,
                 align=PP_ALIGN.RIGHT)
        add_text(slide, Inches(10.5), y + Inches(0.45), Inches(2.0), Inches(0.4),
                 ref, size=12, bold=True, color=color,
                 align=PP_ALIGN.RIGHT)

    # Bottom narrative arc
    by = Inches(7.05)
    add_text(slide, Inches(0.5), by, SW - Inches(1.0), Inches(0.3),
             "→  The reader can interrogate any act by clicking a country, picking a year, or selecting a role bar.",
             size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# =========================================================================
# SLIDE 19, 20 — KEY INSIGHTS (8 in 2 slides)
# =========================================================================
def insight_card(slide, x, y, w, h, num, title, body, color):
    card(slide, x, y, w, h, fill=WHITE, border=color)
    add_rect(slide, x, y, Inches(0.7), h, fill=color)
    add_text(slide, x, y, Inches(0.7), h,
             num, size=28, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.85), y + Inches(0.1), w - Inches(1.0), Inches(0.4),
             title, size=13, bold=True, color=FOREST)
    add_text(slide, x + Inches(0.85), y + Inches(0.5), w - Inches(1.0), h - Inches(0.55),
             body, size=11, color=DARK, line_spacing=1.3)


def slide_insights_1():
    slide = slide_chrome("Key Insights — Part 1", "What the data revealed (1 / 2)", page_n=19)

    insights = [
        ("01", "ML Engineering is the highest-paid AI/ML track",
         "ML Engineer / MLOps leads at $194,303 — 81% above Data Analyst. Software Engineer ($180K) and Research Scientist ($175K) follow.",
         FOREST),
        ("02", "Compensation grows nearly 2× from Entry to Executive",
         "Entry $102,743 → Mid $140,382 → Senior $170,351 → Executive $195,935. The EX/EN ratio is exactly 1.91×.",
         EMERALD),
        ("03", "Remote work has collapsed since 2022",
         "Remote dropped from 53.5% (2022) to 20.3% (2025). On-site nearly doubled from 43.5% to 79.6%. Hybrid effectively vanished.",
         JADE),
        ("04", "The US dominates volume — and the premium is everywhere",
         "United States supplies 90.4% of submissions (133,143 rows) at $161,381 average. Rest of World averages $108,024.",
         HONEY),
    ]
    cols = 2; rows = 2
    cw = Inches(6.0); ch = Inches(2.4); gx = Inches(0.3); gy = Inches(0.25)
    x0 = Inches(0.5); y0 = Inches(1.7)
    for i, (n, t, b, c) in enumerate(insights):
        col = i % cols; row = i // cols
        x = x0 + col * (cw + gx); y = y0 + row * (ch + gy)
        insight_card(slide, x, y, cw, ch, n, t, b, c)


def slide_insights_2():
    slide = slide_chrome("Key Insights — Part 2", "What the data revealed (2 / 2)", page_n=20)

    insights = [
        ("05", "Salaries are concentrated in $100K–$200K",
         "55% of all records sit in the $100K–$200K window (31% in $100K–$150K + 24% in $150K–$200K). Top tier $260K+ holds 7.3%.",
         RUST),
        ("06", "AI/ML Core commands a 13% premium",
         "AI/ML Core averages $171,009 vs $150,685 for Data & Tech. Direct AI/ML exposure earns ~$20K extra per year.",
         FOREST),
        ("07", "Company size matters — but only past Small",
         "Small $103K → Medium $155K → Large $165K. The biggest jump is Small→Medium (+50%); Medium→Large is only $10K.",
         EMERALD),
        ("08", "The market is maturing — volume up, pay flat",
         "Submission volume grew 48× from 2022 (1,587) to 2025 (76,424). Median salary trajectory plateaus after 2024 ($149K → $146K).",
         JADE),
    ]
    cols = 2; rows = 2
    cw = Inches(6.0); ch = Inches(2.4); gx = Inches(0.3); gy = Inches(0.25)
    x0 = Inches(0.5); y0 = Inches(1.7)
    for i, (n, t, b, c) in enumerate(insights):
        col = i % cols; row = i // cols
        x = x0 + col * (cw + gx); y = y0 + row * (ch + gy)
        insight_card(slide, x, y, cw, ch, n, t, b, c)


# =========================================================================
# SLIDE 21, 22 — HOW TABLEAU HELPED
# =========================================================================
def how_card(slide, x, y, w, h, head, body, color):
    card(slide, x, y, w, h, fill=WHITE, border=color)
    add_rect(slide, x, y, w, Inches(0.4), fill=color)
    add_text(slide, x, y, w, Inches(0.4),
             head, size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.2), y + Inches(0.5), w - Inches(0.4), h - Inches(0.55),
             body, size=11, color=DARK, line_spacing=1.35)


def slide_tableau_helped_1():
    slide = slide_chrome("How Tableau Helped (1 / 2)", "From static figures to a tool that stakeholders can interrogate", page_n=21)

    items = [
        ("01  ·  Single Interactive Canvas",
         "The 20 Python figures each told one story. Tableau merged the strongest 10 into a unified 1200×1500 dashboard with cross-filtering — investigate any cohort without producing a new figure.",
         FOREST),
        ("02  ·  Click-to-Filter Actions",
         "A single Filter action on the Global Map propagates the selected country into every other worksheet — synchronised cross-analysis impractical to replicate in matplotlib.",
         EMERALD),
        ("03  ·  Custom Calculated Fields",
         "Avg Salary USD, Median Salary USD, EX/EN Ratio, custom shelf sorts on salary_band and title_group. Consistent metrics across every worksheet.",
         JADE),
        ("04  ·  Geographic Intelligence",
         "Tagging company_location with Country/Region role auto-resolved 87 ISO-2 codes into a world choropleth — no GIS code required.",
         HONEY),
    ]
    cols = 2; rows = 2
    cw = Inches(6.0); ch = Inches(2.4); gx = Inches(0.3); gy = Inches(0.25)
    x0 = Inches(0.5); y0 = Inches(1.7)
    for i, (h, b, c) in enumerate(items):
        col = i % cols; row = i // cols
        x = x0 + col * (cw + gx); y = y0 + row * (ch + gy)
        how_card(slide, x, y, cw, ch, h, b, c)


def slide_tableau_helped_2():
    slide = slide_chrome("How Tableau Helped (2 / 2)", "From explore to communicate — Tableau's compounding advantages", page_n=22)

    items = [
        ("05  ·  Tooltips Beat Static Legends",
         "Every panel has a custom tooltip (count, avg, median, deviation). Hovering replaces the need for legends — every chart self-documents.",
         RUST),
        ("06  ·  KPI Strip for At-A-Glance Reading",
         "Five colour-coded KPI tiles ($147K · $156K · 147,348 · 20.8% · 1.91×) sit under the title bar — readers grasp the market in two seconds.",
         FOREST),
        ("07  ·  Visual Encoding Richness",
         "Sequential color (choropleth · role bar), categorical palettes (Work Mode · Salary Band), dual-pane on Trend (median line + volume bars) — all from one tool.",
         EMERALD),
        ("08  ·  Reproducibility via Packaged .twb",
         "Entire dashboard is a single Dashboard.twb XML file. Anyone with Tableau Desktop can repoint the data and reproduce every chart — no hidden state.",
         JADE),
    ]
    cols = 2; rows = 2
    cw = Inches(6.0); ch = Inches(2.4); gx = Inches(0.3); gy = Inches(0.25)
    x0 = Inches(0.5); y0 = Inches(1.7)
    for i, (h, b, c) in enumerate(items):
        col = i % cols; row = i // cols
        x = x0 + col * (cw + gx); y = y0 + row * (ch + gy)
        how_card(slide, x, y, cw, ch, h, b, c)

    # Bottom callout
    by = Inches(6.7)
    card(slide, Inches(0.5), by, SW - Inches(1.0), Inches(0.55),
         fill=LIGHT_MINT, border=EMERALD)
    add_multi_text(slide, Inches(0.7), by, Inches(12), Inches(0.55), runs=[
        {"text": "✦  In one sentence  ·  ", "size": 13, "bold": True, "color": EMERALD},
        {"text": "Tableau converted our static findings into a tool. Without it the report is a deck of charts; with it, the report has a companion artefact that the reader can manipulate.",
         "size": 11, "italic": True, "color": DARK},
    ], anchor=MSO_ANCHOR.MIDDLE)


# =========================================================================
# SLIDE 23 — DESIGN SYSTEM
# =========================================================================
def slide_design():
    slide = slide_chrome("Design System  ·  Forest Emerald", "Color palette and visual language used across dashboard and report", page_n=23)

    swatches = [
        ("#1B4332", FOREST,  "Forest",        "Primary  ·  title bar, footer, KPI #1"),
        ("#2D6A4F", EMERALD, "Emerald",       "Subtitle, all panel borders (2px)"),
        ("#40916C", JADE,    "Jade",          "KPI #2  ·  Mean salary"),
        ("#52B788", SAGE,    "Sage",          "KPI #3  ·  Total records"),
        ("#D8973C", HONEY,   "Honey",         "KPI #4  ·  Remote workers, insight border"),
        ("#BC4749", RUST,    "Rust",          "KPI #5  ·  EX/EN multiple"),
        ("#FAF8F0", CREAM,   "Cream",         "Card / chart background"),
        ("#E8F1ED", LIGHT_MINT, "Mint Mist",  "Section header bands"),
        ("#FFF4E5", PEACH,   "Warm Peach",    "Insight callout backgrounds"),
    ]
    cols = 3; rows = 3
    cw = Inches(4.05); ch = Inches(1.5); gx = Inches(0.15); gy = Inches(0.15)
    x0 = Inches(0.5); y0 = Inches(1.7)
    for i, (hex_, c, name, role) in enumerate(swatches):
        col = i % cols; row = i // cols
        x = x0 + col * (cw + gx); y = y0 + row * (ch + gy)
        card(slide, x, y, cw, ch, fill=WHITE, border=c)
        # color block
        add_rect(slide, x + Inches(0.1), y + Inches(0.1), Inches(1.3), ch - Inches(0.2), fill=c)
        # name + hex
        add_text(slide, x + Inches(1.55), y + Inches(0.1), cw - Inches(1.7), Inches(0.4),
                 name, size=14, bold=True, color=FOREST)
        add_text(slide, x + Inches(1.55), y + Inches(0.5), cw - Inches(1.7), Inches(0.3),
                 hex_, size=10, color=GRAY, font="Consolas")
        add_text(slide, x + Inches(1.55), y + Inches(0.85), cw - Inches(1.7), Inches(0.55),
                 role, size=10, italic=True, color=DARK)

    # Bottom design notes
    by = Inches(6.5)
    add_text(slide, Inches(0.5), by, SW - Inches(1.0), Inches(0.4),
             "Why this palette? — Distinctive (most dashboards default to blue), accessible "
             "(strong contrast), and earthy (signals the maturity of the AI labour market).",
             size=11, italic=True, color=GRAY, align=PP_ALIGN.CENTER)


# =========================================================================
# SLIDE 24 — TECH STACK
# =========================================================================
def slide_stack():
    slide = slide_chrome("Tech Stack & Deliverables", "What is in the repository", page_n=24)

    items = [
        ("DATA",         "salaries.csv · salaries_clean_final.csv · salaries_enhanced.csv",   FOREST),
        ("PYTHON",       "5 preprocessing scripts · pandas · numpy · scipy",                   EMERALD),
        ("EDA",          "AI_Job_Market_EDA.ipynb · 20 figures · 10 summary CSVs",             JADE),
        ("TABLEAU",      "Dashboard.twb (10 worksheets) · Dashboard.pdf · build plan",         SAGE),
        ("REPORT",       "Final_Project_Report.tex / .pdf  ·  40 pages  ·  9 chapters",       HONEY),
        ("PRESENT",      "Group8_Final_Presentation.pptx (this deck)",                        RUST),
    ]
    rh = Inches(0.65); y0 = Inches(1.7); gap = Inches(0.12)
    for i, (label, body, c) in enumerate(items):
        y = y0 + i * (rh + gap)
        # vertical accent
        add_rect(slide, Inches(0.5), y, Inches(0.15), rh, fill=c)
        # card
        card(slide, Inches(0.7), y, SW - Inches(1.2), rh, fill=WHITE, border=c)
        # label
        add_text(slide, Inches(0.9), y, Inches(2.0), rh,
                 label, size=14, bold=True, color=c, anchor=MSO_ANCHOR.MIDDLE)
        # body
        add_text(slide, Inches(2.9), y, Inches(9.5), rh,
                 body, size=11, color=DARK, anchor=MSO_ANCHOR.MIDDLE,
                 font="Consolas")


# =========================================================================
# SLIDE 25 — CONCLUSION
# =========================================================================
def slide_conclusion():
    slide = slide_chrome("Conclusion & Takeaways", "What we built and why it matters", page_n=25)

    # 4 takeaway cards
    items = [
        ("END-TO-END",
         "From raw 151,445 rows to an interactive 10-worksheet Tableau dashboard, with a Python pipeline and a 40-page report tying it all together.",
         FOREST),
        ("EVIDENCE-BASED",
         "Every claim in the dashboard is backed by a real CC0 record. No synthetic data introduced anywhere in the pipeline.",
         EMERALD),
        ("INTERACTIVE",
         "Click any country to re-scope every panel. KPI strip gives a 2-second read. Hover for tooltip-level detail.",
         JADE),
        ("REPRODUCIBLE",
         "Single .twb file + 5 Python scripts + a CSV. Anyone can clone, repoint, and regenerate.",
         HONEY),
    ]
    cols = 2; rows = 2
    cw = Inches(6.0); ch = Inches(1.85); gx = Inches(0.3); gy = Inches(0.25)
    x0 = Inches(0.5); y0 = Inches(1.7)
    for i, (h, b, c) in enumerate(items):
        col = i % cols; row = i // cols
        x = x0 + col * (cw + gx); y = y0 + row * (ch + gy)
        card(slide, x, y, cw, ch, fill=WHITE, border=c)
        add_rect(slide, x, y, cw, Inches(0.45), fill=c)
        add_text(slide, x, y, cw, Inches(0.45),
                 h, size=14, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.2), y + Inches(0.55), cw - Inches(0.4), ch - Inches(0.6),
                 b, size=11, color=DARK, line_spacing=1.4)

    # Final line
    by = Inches(6.0)
    card(slide, Inches(0.5), by, SW - Inches(1.0), Inches(1.1), fill=LIGHT_MINT, border=EMERALD)
    add_multi_text(slide, Inches(0.7), by + Inches(0.1), Inches(12.0), Inches(0.4), runs=[
        {"text": "FIVE GUIDING QUESTIONS  ·  ALL ANSWERED:",
         "size": 13, "bold": True, "color": EMERALD},
    ])
    add_multi_text(slide, Inches(0.7), by + Inches(0.45), Inches(12.0), Inches(0.7), runs=[
        {"text": "ML Engineer at $194K  ·  ", "size": 11, "bold": True, "color": FOREST},
        {"text": "1.91× from Entry to Exec  ·  ", "size": 11, "bold": True, "color": FOREST},
        {"text": "Remote 53% → 20%  ·  ",       "size": 11, "bold": True, "color": FOREST},
        {"text": "US 90% of submissions  ·  ",  "size": 11, "bold": True, "color": FOREST},
        {"text": "55% in $100–$200K band",      "size": 11, "bold": True, "color": FOREST},
    ])


# =========================================================================
# SLIDE 26 — THANK YOU
# =========================================================================
def slide_thank_you():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, SW, SH, fill=FOREST)

    # Decorative accents
    add_rect(slide, Inches(-1), Inches(5.5), Inches(8), Inches(2.5), fill=EMERALD)
    add_rect(slide, Inches(8), Inches(0), Inches(8), Inches(1.2), fill=EMERALD)

    # Triangle accent
    tri = slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                 Inches(11.0), Inches(5.5), Inches(2.5), Inches(2.0))
    tri.shadow.inherit = False
    tri.fill.solid(); tri.fill.fore_color.rgb = SAGE
    tri.line.fill.background()

    # Sage horizontal accent
    add_rect(slide, Inches(0.7), Inches(2.8), Inches(1.5), Inches(0.06), fill=SAGE)

    # Big text
    add_text(slide, Inches(0.7), Inches(2.2), Inches(12), Inches(0.6),
             "Q  &  A", size=18, bold=True, color=SAGE)
    add_text(slide, Inches(0.7), Inches(3.0), Inches(12), Inches(1.5),
             "THANK YOU", size=88, bold=True, color=WHITE)
    add_text(slide, Inches(0.7), Inches(4.5), Inches(12), Inches(0.6),
             "We welcome your questions on data, methodology, or the dashboard.",
             size=18, italic=True, color=CREAM)

    # Footer info
    add_text(slide, Inches(0.7), Inches(6.0), Inches(12), Inches(0.4),
             "GROUP 8  ·  CS3012  ·  FAST-NUCES Islamabad  ·  Spring 2026",
             size=12, bold=True, color=SAGE)
    add_text(slide, Inches(0.7), Inches(6.45), Inches(12), Inches(0.4),
             "Muhammad Nouman Hafeez (21i-0416)  ·  Muhammad Asim (21i-0852)  ·  Sara Jabeen (21i-0624)",
             size=11, color=CREAM)
    add_text(slide, Inches(0.7), Inches(6.8), Inches(12), Inches(0.4),
             "Repository: github.com/noumanic/ai-job-market-analytics  ·  Source: aijobs.net (CC0)",
             size=11, italic=True, color=CREAM)


# =========================================================================
# Apply transitions to every slide
# =========================================================================
def apply_transitions(prs):
    """Inject a fade transition on every slide via raw XML."""
    transition_xml = (
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'spd="med"><p:fade/></p:transition>'
    )
    for slide in prs.slides:
        sld = slide._element
        # Insert <p:transition/> before <p:timing/> if present, else at the end
        new = etree.fromstring(transition_xml)
        sld.append(new)


# =========================================================================
# BUILD
# =========================================================================
def build():
    slide_title()
    slide_agenda()
    slide_group()
    slide_overview()
    slide_questions()
    slide_dataset()
    slide_authenticity()
    slide_pipeline_overview()
    slide_stage_clean()
    slide_stage_grouping()
    slide_stages_4_5()
    slide_eda_overview()
    slide_tableau_master()
    slide_sheets_group_1()
    slide_sheets_group_2()
    slide_sheets_group_3()
    slide_sheets_group_4()
    slide_story()
    slide_insights_1()
    slide_insights_2()
    slide_tableau_helped_1()
    slide_tableau_helped_2()
    slide_design()
    slide_stack()
    slide_conclusion()
    slide_thank_you()

    apply_transitions(prs)
    prs.save(str(OUT))
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
